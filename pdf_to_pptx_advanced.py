# -*- coding: utf-8 -*-
"""
高度なPDF→PPTX変換ツール
NotebookLMで生成されたスライドPDFを含む、様々なPDFに対応

3つのモード:
1. extract: PDF内のテキスト/画像を直接抽出（通常のPDF向け）
2. ocr: OCRでテキストを認識（画像ベースのPDF向け、NotebookLM等）
3. hybrid: 両方を組み合わせ（自動判断）

使い方:
    python pdf_to_pptx_advanced.py input.pdf output.pptx --mode hybrid
"""

import fitz  # PyMuPDF
import pdfplumber
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import os
import sys
import argparse
import re
import tempfile
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from pathlib import Path

# 文字化け対策
try:
    import ftfy
    HAS_FTFY = True
except ImportError:
    HAS_FTFY = False

# OCR
try:
    import pytesseract
    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False


# =============================================================================
# データクラス定義
# =============================================================================

@dataclass
class TextBlock:
    """テキストブロック"""
    text: str
    x: float
    y: float
    width: float
    height: float
    font_size: float = 12.0
    is_bold: bool = False
    color: Tuple[int, int, int] = (0, 0, 0)
    block_type: str = "text"
    confidence: float = 1.0  # OCRの信頼度


@dataclass
class ImageBlock:
    """画像ブロック"""
    image_data: bytes
    x: float
    y: float
    width: float
    height: float


@dataclass
class PageElements:
    """ページの全要素"""
    page_num: int
    width: float
    height: float
    texts: List[TextBlock] = field(default_factory=list)
    images: List[ImageBlock] = field(default_factory=list)
    page_image: Optional[bytes] = None  # ページ全体の画像


# =============================================================================
# 文字化け修正
# =============================================================================

MOJIBAKE_MAP = {
    "縺": "", "繧": "", "繝": "",
    "笆": "★", "窶": "ー",
}

def fix_mojibake(text: str) -> str:
    """文字化けを修正"""
    if not text:
        return text

    if HAS_FTFY:
        text = ftfy.fix_text(text)

    for wrong, correct in MOJIBAKE_MAP.items():
        text = text.replace(wrong, correct)

    text = ''.join(char for char in text if ord(char) >= 32 or char in '\n\t')
    return text.strip()


def is_valid_japanese_text(text: str) -> bool:
    """有効な日本語テキストかチェック"""
    if not text or len(text) < 2:
        return False

    # 文字化けパターン
    mojibake_chars = sum(1 for c in text if c in '縺繧繝')
    if mojibake_chars > len(text) * 0.2:
        return False

    # 有効な文字（日本語、英数字、記号）の割合
    valid_chars = sum(1 for c in text if (
        '\u3040' <= c <= '\u309F' or  # ひらがな
        '\u30A0' <= c <= '\u30FF' or  # カタカナ
        '\u4E00' <= c <= '\u9FFF' or  # 漢字
        'a' <= c.lower() <= 'z' or
        '0' <= c <= '9' or
        c in ' 　・、。！？「」『』【】（）()[]{}:：;；.,。、-ー〜/／\n\t'
    ))

    return valid_chars > len(text) * 0.5


# =============================================================================
# OCR処理
# =============================================================================

class OCRProcessor:
    """OCR処理エンジン"""

    def __init__(self, lang: str = "jpn+eng"):
        # Tesseractのパスを設定（Windows用）
        if sys.platform == "win32":
            tesseract_paths = [
                r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            ]
            for path in tesseract_paths:
                if os.path.exists(path):
                    if HAS_TESSERACT:
                        pytesseract.pytesseract.tesseract_cmd = path
                    break

        # カスタムtessdataパスを設定（ユーザーホームにjpn.traineddataがある場合）
        self.custom_tessdata = None
        home_tessdata = os.path.expanduser("~")
        if os.path.exists(os.path.join(home_tessdata, "jpn.traineddata")):
            self.custom_tessdata = home_tessdata

        # 利用可能な言語を確認
        self.lang = lang
        if HAS_TESSERACT:
            try:
                available = pytesseract.get_languages()
                if "jpn" not in available and self.custom_tessdata:
                    # カスタムパスに日本語データがある
                    self.lang = "jpn+eng"
                elif "jpn" not in available:
                    # 日本語がなければ英語のみ
                    self.lang = "eng"
            except Exception:
                if self.custom_tessdata:
                    self.lang = "jpn+eng"
                else:
                    self.lang = "eng"

        if not HAS_TESSERACT:
            print("警告: pytesseractがインストールされていません")
            print("  pip install pytesseract")
            print("  + Tesseractのインストールが必要です")

    def process_image(self, image: Image.Image) -> List[TextBlock]:
        """画像からテキストブロックを抽出"""
        if not HAS_TESSERACT:
            return []

        try:
            # 画像サイズ取得
            img_width, img_height = image.size

            # Tesseract OCRでテキスト+位置情報を取得
            config = ""
            if self.custom_tessdata:
                config = f'--tessdata-dir "{self.custom_tessdata}"'

            data = pytesseract.image_to_data(
                image,
                lang=self.lang,
                output_type=pytesseract.Output.DICT,
                config=config
            )

            blocks = []
            current_block = None
            current_text = []

            n_boxes = len(data['text'])
            for i in range(n_boxes):
                text = data['text'][i].strip()
                conf = float(data['conf'][i])

                # 信頼度が低いか空のテキストはスキップ
                if conf < 30 or not text:
                    # ブロックの区切り
                    if current_block and current_text:
                        current_block.text = ' '.join(current_text)
                        current_block.text = fix_mojibake(current_block.text)
                        if is_valid_japanese_text(current_block.text):
                            blocks.append(current_block)
                        current_block = None
                        current_text = []
                    continue

                x = data['left'][i]
                y = data['top'][i]
                w = data['width'][i]
                h = data['height'][i]

                # 新しいブロック開始または既存ブロックに追加
                if current_block is None:
                    current_block = TextBlock(
                        text="",
                        x=x,
                        y=y,
                        width=w,
                        height=h,
                        confidence=conf / 100
                    )
                    current_text = [text]
                else:
                    # Y座標が近い場合は同じ行として扱う
                    if abs(y - current_block.y) < current_block.height * 1.5:
                        current_text.append(text)
                        # ブロックの境界を更新
                        new_right = max(current_block.x + current_block.width, x + w)
                        current_block.width = new_right - current_block.x
                        current_block.height = max(current_block.height, h)
                    else:
                        # 新しい行 = ブロック確定
                        current_block.text = ' '.join(current_text)
                        current_block.text = fix_mojibake(current_block.text)
                        if is_valid_japanese_text(current_block.text):
                            blocks.append(current_block)

                        current_block = TextBlock(
                            text="",
                            x=x,
                            y=y,
                            width=w,
                            height=h,
                            confidence=conf / 100
                        )
                        current_text = [text]

            # 最後のブロック
            if current_block and current_text:
                current_block.text = ' '.join(current_text)
                current_block.text = fix_mojibake(current_block.text)
                if is_valid_japanese_text(current_block.text):
                    blocks.append(current_block)

            # フォントサイズの推定（ブロックの高さから）
            for block in blocks:
                block.font_size = min(72, max(8, block.height * 0.7))

                # 見出しの判定（大きいフォント、上部に位置）
                if block.font_size >= 20 or block.y < img_height * 0.2:
                    block.block_type = "heading"
                    block.is_bold = True
                elif block.text.startswith(('•', '・', '-', '●', '○')):
                    block.block_type = "bullet"

            return blocks

        except Exception as e:
            print(f"    OCRエラー: {e}")
            return []


# =============================================================================
# PDF抽出エンジン
# =============================================================================

class PDFExtractor:
    """PDF要素抽出エンジン"""

    def __init__(self, pdf_path: str, mode: str = "hybrid", verbose: bool = True):
        self.pdf_path = pdf_path
        self.mode = mode
        self.verbose = verbose
        self.doc = fitz.open(pdf_path)
        self.ocr = OCRProcessor() if mode in ("ocr", "hybrid") else None

    def log(self, msg: str):
        if self.verbose:
            print(msg)

    def extract_all(self) -> List[PageElements]:
        """全ページを抽出"""
        pages = []
        total = len(self.doc)

        for i in range(total):
            self.log(f"  ページ {i+1}/{total} を処理中...")
            pages.append(self.extract_page(i))

        return pages

    def extract_page(self, page_num: int) -> PageElements:
        """1ページを抽出"""
        page = self.doc[page_num]
        rect = page.rect

        elements = PageElements(
            page_num=page_num,
            width=rect.width,
            height=rect.height
        )

        # テキスト抽出（extractモードまたはhybrid）
        if self.mode in ("extract", "hybrid"):
            elements.texts = self._extract_texts_direct(page)

        # OCR処理（ocrモードまたはhybridでテキストが少ない場合）
        if self.mode == "ocr" or (self.mode == "hybrid" and len(elements.texts) < 3):
            ocr_texts = self._extract_texts_ocr(page)
            if len(ocr_texts) > len(elements.texts):
                elements.texts = ocr_texts

        # 画像抽出
        elements.images = self._extract_images(page)

        # ページ全体の画像（背景用）
        elements.page_image = self._render_page_image(page)

        return elements

    def _extract_texts_direct(self, page: fitz.Page) -> List[TextBlock]:
        """PyMuPDFで直接テキスト抽出"""
        blocks = []
        text_dict = page.get_text("dict")

        for block in text_dict.get("blocks", []):
            if block.get("type") != 0:
                continue

            text = ""
            font_size = 12.0
            is_bold = False
            color = (0, 0, 0)

            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text += span.get("text", "")
                    if not font_size or font_size == 12.0:
                        font_size = span.get("size", 12.0)
                        flags = span.get("flags", 0)
                        is_bold = bool(flags & 16)
                        c = span.get("color", 0)
                        if isinstance(c, int):
                            color = ((c >> 16) & 0xFF, (c >> 8) & 0xFF, c & 0xFF)
                text += "\n"

            text = fix_mojibake(text.strip())
            if not text or not is_valid_japanese_text(text):
                continue

            bbox = block.get("bbox", (0, 0, 100, 20))
            block_type = "heading" if font_size >= 20 else "text"

            blocks.append(TextBlock(
                text=text,
                x=bbox[0],
                y=bbox[1],
                width=bbox[2] - bbox[0],
                height=bbox[3] - bbox[1],
                font_size=font_size,
                is_bold=is_bold,
                color=color,
                block_type=block_type
            ))

        return blocks

    def _extract_texts_ocr(self, page: fitz.Page) -> List[TextBlock]:
        """OCRでテキスト抽出"""
        if not self.ocr:
            return []

        # ページを画像化
        mat = fitz.Matrix(2.0, 2.0)  # 2倍解像度
        pix = page.get_pixmap(matrix=mat)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

        # OCR実行
        ocr_blocks = self.ocr.process_image(img)

        # 座標をPDFスケールに変換
        scale = 0.5  # 2倍解像度の逆
        for block in ocr_blocks:
            block.x *= scale
            block.y *= scale
            block.width *= scale
            block.height *= scale

        return ocr_blocks

    def _extract_images(self, page: fitz.Page) -> List[ImageBlock]:
        """画像を抽出"""
        images = []

        for img_info in page.get_images(full=True):
            try:
                xref = img_info[0]
                base_img = self.doc.extract_image(xref)
                img_data = base_img.get("image")

                if not img_data:
                    continue

                rects = page.get_image_rects(xref)
                if rects:
                    r = rects[0]
                    images.append(ImageBlock(
                        image_data=img_data,
                        x=r.x0,
                        y=r.y0,
                        width=r.width,
                        height=r.height
                    ))
            except Exception as e:
                self.log(f"    画像抽出エラー: {e}")

        return images

    def _render_page_image(self, page: fitz.Page) -> bytes:
        """ページ全体を画像としてレンダリング"""
        mat = fitz.Matrix(1.5, 1.5)  # 1.5倍解像度
        pix = page.get_pixmap(matrix=mat)
        return pix.tobytes("png")

    def close(self):
        self.doc.close()


# =============================================================================
# PPTX生成
# =============================================================================

class PPTXBuilder:
    """PPTX生成エンジン"""

    PT_TO_EMU = 914400 / 72

    def __init__(self, output_path: str, use_background: bool = True, verbose: bool = True):
        self.output_path = output_path
        self.use_background = use_background
        self.verbose = verbose

        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def log(self, msg: str):
        if self.verbose:
            print(msg)

    def build(self, pages: List[PageElements]):
        """PPTXを生成"""
        blank = self.prs.slide_layouts[6]

        for page in pages:
            self.log(f"  スライド {page.page_num + 1}/{len(pages)} を生成中...")

            slide = self.prs.slides.add_slide(blank)

            # スケール計算
            scale_x = self.prs.slide_width / (page.width * self.PT_TO_EMU)
            scale_y = self.prs.slide_height / (page.height * self.PT_TO_EMU)
            scale = min(scale_x, scale_y)

            offset_x = (self.prs.slide_width - page.width * self.PT_TO_EMU * scale) / 2
            offset_y = (self.prs.slide_height - page.height * self.PT_TO_EMU * scale) / 2

            # 背景画像
            if self.use_background and page.page_image:
                self._add_background(slide, page.page_image)

            # 画像（背景でない場合）
            if not self.use_background:
                for img in page.images:
                    self._add_image(slide, img, scale, offset_x, offset_y)

            # テキスト
            for text in page.texts:
                self._add_text(slide, text, scale, offset_x, offset_y)

        self.prs.save(self.output_path)

    def _add_background(self, slide, image_data: bytes):
        """背景画像を追加"""
        try:
            stream = io.BytesIO(image_data)
            slide.shapes.add_picture(
                stream, 0, 0,
                width=self.prs.slide_width,
                height=self.prs.slide_height
            )
        except Exception as e:
            self.log(f"    背景画像エラー: {e}")

    def _add_image(self, slide, img: ImageBlock, scale: float, ox: float, oy: float):
        """画像を追加"""
        try:
            px = int(img.x * self.PT_TO_EMU * scale + ox)
            py = int(img.y * self.PT_TO_EMU * scale + oy)
            pw = int(img.width * self.PT_TO_EMU * scale)
            ph = int(img.height * self.PT_TO_EMU * scale)

            stream = io.BytesIO(img.image_data)
            slide.shapes.add_picture(stream, px, py, pw, ph)
        except Exception as e:
            self.log(f"    画像追加エラー: {e}")

    def _add_text(self, slide, text: TextBlock, scale: float, ox: float, oy: float):
        """テキストボックスを追加"""
        px = int(text.x * self.PT_TO_EMU * scale + ox)
        py = int(text.y * self.PT_TO_EMU * scale + oy)
        pw = max(int(text.width * self.PT_TO_EMU * scale), Inches(1))
        ph = max(int(text.height * self.PT_TO_EMU * scale), Inches(0.3))

        box = slide.shapes.add_textbox(px, py, pw, ph)
        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text.text

        # フォントサイズ（スケール適用、範囲制限）
        size = text.font_size * scale
        p.font.size = Pt(max(8, min(48, size)))
        p.font.bold = text.is_bold

        r, g, b = text.color
        p.font.color.rgb = RGBColor(r, g, b)


# =============================================================================
# メイン処理
# =============================================================================

def convert(input_pdf: str, output_pptx: str, mode: str = "hybrid",
            use_background: bool = True, verbose: bool = True):
    """PDFをPPTXに変換"""

    if verbose:
        print("=" * 60)
        print("PDF to PPTX Converter (Advanced)")
        print("=" * 60)
        print(f"入力: {input_pdf}")
        print(f"出力: {output_pptx}")
        print(f"モード: {mode}")
        print(f"背景画像: {'あり' if use_background else 'なし'}")
        print()

    if not os.path.exists(input_pdf):
        print(f"エラー: ファイルが見つかりません: {input_pdf}")
        sys.exit(1)

    # 抽出
    if verbose:
        print("PDF解析中...")

    extractor = PDFExtractor(input_pdf, mode=mode, verbose=verbose)
    pages = extractor.extract_all()
    extractor.close()

    # 統計
    if verbose:
        total_texts = sum(len(p.texts) for p in pages)
        total_images = sum(len(p.images) for p in pages)
        print()
        print(f"抽出結果:")
        print(f"  ページ数: {len(pages)}")
        print(f"  テキストブロック: {total_texts}")
        print(f"  画像: {total_images}")
        print()

    # 生成
    if verbose:
        print("PPTX生成中...")

    builder = PPTXBuilder(output_pptx, use_background=use_background, verbose=verbose)
    builder.build(pages)

    if verbose:
        size = os.path.getsize(output_pptx) / 1024
        print()
        print("=" * 60)
        print("変換完了!")
        print(f"  出力: {output_pptx}")
        print(f"  サイズ: {size:.1f} KB")
        print("=" * 60)


def main():
    parser = argparse.ArgumentParser(
        description="PDFを編集可能なPPTXに変換（NotebookLM対応）",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
モード説明:
  extract  - PDF内のテキストを直接抽出（通常のPDF向け）
  ocr      - OCRでテキスト認識（画像ベースのPDF向け）
  hybrid   - 両方を組み合わせ（自動判断、推奨）

使用例:
  # 推奨（自動判断）
  python pdf_to_pptx_advanced.py slide.pdf output.pptx

  # NotebookLMのPDF（OCRモード）
  python pdf_to_pptx_advanced.py slide.pdf output.pptx --mode ocr

  # 背景画像なし（テキストのみ抽出）
  python pdf_to_pptx_advanced.py slide.pdf output.pptx --no-background

必要条件:
  - PyMuPDF, pdfplumber, python-pptx
  - OCRモード: pytesseract + Tesseract OCR
        """
    )

    parser.add_argument("input", help="入力PDFファイル")
    parser.add_argument("output", help="出力PPTXファイル")
    parser.add_argument("--mode", "-m", choices=["extract", "ocr", "hybrid"],
                       default="hybrid", help="抽出モード（デフォルト: hybrid）")
    parser.add_argument("--no-background", action="store_true",
                       help="背景画像を含めない")
    parser.add_argument("--quiet", "-q", action="store_true",
                       help="詳細ログを抑制")

    args = parser.parse_args()

    convert(
        args.input,
        args.output,
        mode=args.mode,
        use_background=not args.no_background,
        verbose=not args.quiet
    )


if __name__ == "__main__":
    main()
